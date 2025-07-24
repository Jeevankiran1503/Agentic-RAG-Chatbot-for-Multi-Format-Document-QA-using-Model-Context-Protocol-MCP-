import os
import pandas as pd
import docx
import fitz  
from pptx import Presentation

SUPPORTED_EXTENSIONS = [".txt", ".md", ".pdf", ".docx", ".pptx", ".csv"]

def load_documents(path):
    docs = []

    if os.path.isdir(path):
        for file in os.listdir(path):
            full_path = os.path.join(path, file)
            docs.extend(load_documents(full_path))  
    elif os.path.isfile(path):
        ext = os.path.splitext(path)[1].lower()

        if ext in [".txt", ".md"]:
            with open(path, "r", encoding="utf-8") as f:
                docs.append(f.read())

        elif ext == ".pdf":
            text = ""
            with fitz.open(path) as doc:
                for page in doc:
                    text += page.get_text()
            docs.append(text)

        elif ext == ".docx":
            doc = docx.Document(path)
            text = "\n".join([para.text for para in doc.paragraphs])
            docs.append(text)

        elif ext == ".pptx":
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            docs.append(text)

        elif ext == ".csv":
            try:
                df = pd.read_csv(path, encoding="utf-8")
            except UnicodeDecodeError:
                df = pd.read_csv(path, encoding="latin1")  # fallback for Windows encoding
            text = df.to_string(index=False)
            docs.append(text)

    return docs

def basic_chunk_by_paragraph(text, min_chunk_length=50):
    paragraphs = text.strip().split("\n\n")
    chunks = [para.strip() for para in paragraphs if len(para.strip()) >= min_chunk_length]
    return chunks
